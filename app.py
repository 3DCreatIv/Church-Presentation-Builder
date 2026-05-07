import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO

# App UI
st.set_page_config(page_title="Church Presentation Builder", layout="wide")
st.title("⛪ Church Media Builder")
st.subheader("Upload media and add lyrics to create your service slides")

# Sidebar - Presentation Settings
with st.sidebar:
    st.header("Settings")
    theme_color = st.color_picker("Slide Text Color", "#FFFFFF")
    bg_color = st.color_picker("Background Color", "#000000")
    font_size = st.slider("Font Size", 24, 80, 44)

# Multi-file Uploader
uploaded_files = st.file_uploader("Drag & Drop Images/Videos", accept_multiple_files=True, type=['png', 'jpg', 'jpeg', 'mp4'])
scripture_text = st.text_area("Paste Lyrics or Scripture here (One slide per paragraph)")

def generate_pptx(files, lyrics):
    prs = Presentation()
    
    # 1. Add Lyrics Slides
    if lyrics:
        paragraphs = lyrics.split('\n\n')
        for para in paragraphs:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            # Set Background
            fill = slide.background.fill
            fill.solid()
            # Note: Background color setting requires complex XML in python-pptx, 
            # so we focus on text placement here.
            
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = txBox.text_frame
            tf.text = para
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.size = Pt(font_size)

    # 2. Add Media Slides
    for file in files:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if file.type.startswith('image'):
            # Save temp image to add to slide
            slide.shapes.add_picture(file, Inches(0.5), Inches(0.5), height=Inches(6.5))
        elif file.type.startswith('video'):
            st.info(f"Video detected: {file.name}. Note: Videos are best embedded manually in PPT.")

    # Save to memory
    binary_output = BytesIO()
    prs.save(binary_output)
    return binary_output.getvalue()

# Action Button
if st.button("Generate Presentation"):
    if not uploaded_files and not scripture_text:
        st.warning("Please add some content first!")
    else:
        ppt_data = generate_pptx(uploaded_files, scripture_text)
        st.download_button(
            label="📥 Download PowerPoint",
            data=ppt_data,
            file_name="church_service.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
import streamlit as st
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO

st.set_page_config(page_title="Church Media Builder", layout="wide")
st.title("⛪ Church Media Builder (KJV Edition)")

# 1. API Fetcher for KJV Scripture
def fetch_kjv_verse(reference):
    url = f"https://bible-api.com/

{reference}?translation=kjv"
    try:
        response = requests.get(url)
        if response.status_code == 200:
            data = response.json()
            return f"{data['text']}\n— {data['reference']}"
        else:
            return "Error: Could not find that verse. Check your spelling (e.g., 'John 3:16')."
    except Exception as e:
        return f"Connection error: {e}"

# Sidebar & Layout
with st.sidebar:
    st.header("Settings")
    font_size = st.slider("Font Size", 24, 80, 44)
    
# Layout Columns
col1, col2 = st.columns([1, 1])

with col1:
    st.subheader("1. Add Media")
    uploaded_files = st.file_uploader("Upload Images/Videos", accept_multiple_files=True, type=['png', 'jpg', 'jpeg', 'mp4'])

with col2:
    st.subheader("2. Fetch KJV Scripture")
    verse_ref = st.text_input("Enter Reference (e.g., Psalm 23:1)", placeholder="John 3:16")
    if st.button("Fetch & Add to Lyrics"):
        result = fetch_kjv_verse(verse_ref)
        # Append to existing text area
        if "bible_text" not in st.session_state:
            st.session_state.bible_text = ""
        st.session_state.bible_text += f"\n\n{result}"

# Text Area for Lyrics/Scripture (Persisted in Session State)
if "bible_text" not in st.session_state:
    st.session_state.bible_text = ""

all_lyrics = st.text_area("Final Slide Text (Edit as needed)", value=st.session_state.bible_text, height=200)

# Presentation Generation Logic
def generate_pptx(files, lyrics):
    prs = Presentation()
    
    # Process Lyrics/Scripture
    if lyrics:
        # Splits text by double newlines to create separate slides
        paragraphs = [p.strip() for p in lyrics.split('\n\n') if p.strip()]
        for para in paragraphs:
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank slide
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = para
            p.font.size = Pt(font_size)
            p.alignment = PP_ALIGN.CENTER

    # Add Images
    for file in files:
        if file.type.startswith('image'):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(file, Inches(0.5), Inches(0.5), height=Inches(6.5))
            
    binary_output = BytesIO()
    prs.save(binary_output)
    return binary_output.getvalue()

if st.button("Generate & Download PPT"):
    ppt_data = generate_pptx(uploaded_files, all_lyrics)
    st.download_button("📥 Download KJV Presentation", ppt_data, "service.pptx")
import streamlit as st
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO

st.set_page_config(page_title="KJV Church Media Builder", layout="wide")
st.title("⛪ KJV Church Media Builder")

# Sidebar - Global Settings
with st.sidebar:
    st.header("Slide Settings")
    layout_mode = st.radio("Select Layout Mode", ["Full Screen", "Lower Third"])
    font_size = st.slider("Font Size", 20, 80, 44 if layout_mode == "Full Screen" else 32)
    text_color = st.color_picker("Text Color", "#FFFFFF")
    
    st.info("💡 'Lower Third' mode puts text at the bottom for live stream overlays.")

# Column 1: Scripture & Lyrics
col1, col2 = st.columns([2, 1])

with col2:
    st.subheader("Fetch KJV Scripture")
    verse_ref = st.text_input("Enter Reference", placeholder="John 3:16")
    if st.button("Add Verse to List"):
        res = requests.get(f"https://bible-api.com{verse_ref}?translation=kjv")
        if res.status_code == 200:
            data = res.json()
            verse_text = f"{data['text'].strip()}\n({data['reference']})"
            if "content_list" not in st.session_state: st.session_state.content_list = ""
            st.session_state.content_list += f"\n\n{verse_text}"
        else:
            st.error("Verse not found!")

with col1:
    st.subheader("Edit Slides Content")
    content = st.text_area("One paragraph = One slide", 
                          value=st.session_state.get("content_list", ""), 
                          height=300)

uploaded_files = st.file_uploader("Upload Backgrounds/Media", accept_multiple_files=True)

def generate_pptx(text_content, files, mode):
    prs = Presentation()
    # Standard 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Process Text Slides
    paragraphs = [p.strip() for p in text_content.split('\n\n') if p.strip()]
    for para in paragraphs:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        
        # Determine Box Position
        if mode == "Lower Third":
            # Position at bottom: left, top, width, height
            left, top, width, height = Inches(1), Inches(5), Inches(11.3), Inches(2)
        else:
            # Centered Full Screen
            left, top, width, height = Inches(1), Inches(1.5), Inches(11.3), Inches(4.5)
            
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = para
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor.from_string(text_color.replace("#", ""))

    # Process Image Slides
    for file in files:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(file, Inches(0), Inches(0), width=prs.slide_width)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()

if st.button("✨ Create Presentation"):
    ppt = generate_pptx(content, uploaded_files, layout_mode)
    st.download_button("📥 Download PPTX", ppt, "Church_Service.pptx")

def fetch_kjv_verse(reference):
    url = f"
def fetch_kjv_verse(reference):
    url = f"
https://bible-api.com/
{reference}?translation=kjv"
    try:
        response = requests.get(url)
        if response.status_code == 200:
            data = response.json()
            # Loop through individual verse objects from the API
            formatted_verses = []
            for v in data['verses']:
                # Format as "1 In the beginning..."
                formatted_verses.append(f"{v['verse']} {v['text'].strip()}")
            
            # Combine them into one block for the editor
            full_text = " ".join(formatted_verses)
            return f"{full_text}\n— {data['reference']}"
        else:
            return "Error: Verse not found."
    except Exception as e:
        return f"Error: {e}"

{reference}?translation=kjv"
    try:
        response = requests.get(url)
        if response.status_code == 200:
            data = response.json()
            # Loop through individual verse objects from the API
            formatted_verses = []
            for v in data['verses']:
                # Format as "1 In the beginning..."
                formatted_verses.append(f"{v['verse']} {v['text'].strip()}")
            
            # Combine them into one block for the editor
            full_text = " ".join(formatted_verses)
            return f"{full_text}\n— {data['reference']}"
        else:
            return "Error: Verse not found."
    except Exception as e:
        return f"Error: {e}"
